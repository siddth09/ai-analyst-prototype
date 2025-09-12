import streamlit as st
from PyPDF2 import PdfReader
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from analyst import analyze_startup

st.title("🚀 AI Analyst Prototype")

uploaded_file = st.file_uploader("Upload Pitch Deck (PDF/PPTX)", type=["pdf", "pptx"])

if uploaded_file:
    text = ""
    if uploaded_file.name.endswith(".pdf"):
        reader = PdfReader(uploaded_file)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    elif uploaded_file.name.endswith(".pptx"):
        prs = Presentation(uploaded_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    st.subheader("📑 Extracted Content")
    st.write(text[:1000] + "...")  # preview first 1000 chars

    # Encode with embeddings (demo)
    model = SentenceTransformer("all-MiniLM-L6-v2")
    embedding = model.encode([text])[0]

    # Run analysis
    benchmark, risks, recommendation = analyze_startup(text)

    st.subheader("📊 Benchmarking")
    for k, v in benchmark.items():
        st.write(f"- {k}: {v}")

    st.subheader("⚠️ Risk Flags")
    if risks:
        for r in risks:
            st.write(f"- {r}")
    else:
        st.write("No major risk flags detected")

    st.subheader("💡 Recommendation")
    st.write(recommendation)

    st.success("✅ Prototype complete: investor-ready insights generated!")
